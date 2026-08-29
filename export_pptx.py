import os
import glob
from pptx import Presentation
from pptx.util import Inches

def create_pptx_from_slides(exports_dir="exports", output_pptx="SHIVNETRA47_Master_Slide_Deck.pptx"):
    prs = Presentation()
    
    # Set 16:9 widescreen dimensions (13.333 x 7.5 inches = 1920x1080)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    png_files = sorted(glob.glob(os.path.join(exports_dir, "page_*.png")))
    if not png_files:
        print(f"No PNG files found in {exports_dir}")
        return
    
    for f in png_files:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(f, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
    prs.save(output_pptx)
    print(f"Successfully generated {output_pptx} with {len(png_files)} slides ({os.path.getsize(output_pptx)/(1024*1024):.2f} MB)")

if __name__ == "__main__":
    create_pptx_from_slides()
